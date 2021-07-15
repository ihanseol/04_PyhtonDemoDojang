import openpyxl
import win32com.client
from googletrans import Translator
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# PPT Shape 번역 sub 함수 (변경할 Shape, 바꿀 언어, translator 개체) 
def Trans(shape, language, trans):
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        try:
            font = paragraph.runs[0].font
            Fname = font.name
            Fsize = font.size
            Fbold = font.bold
        except:
            pass
        cur_text = paragraph.text
        new_text = trans.translate(cur_text, dest=language).text
        paragraph.text = new_text
        for run in paragraph.runs:
            font = run.font
            font.name = Fname
            font.size = Fsize
            font.bold = Fbold

# PPT 번역 함수, Shape 형태에 따라 구분 동작 (파일명, 작업대상폴더, 결과저장폴더, 바꿀언어)
def GTrans(file,TPath,RPath,language):
    trans=Translator()
    prs=Presentation(TPath + "/" + file)
    Slides = prs.slides
    for slide in Slides:
        Shapes = slide.shapes
        for shape in Shapes:
            #shape에 text_frame이 있는 경우 동작
            if shape.has_text_frame:
                Trans(shape, language,trans)
            # shape가 표일 경우 동작
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        Trans(cell, language,trans)
            # shape가 그룹으로 묶여 있을 경우 동작
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for shp in shape.shapes:
                    if shp.has_text_frame:
                        Trans(shp, language,trans)
    prs.save(RPath + "/" + file)

# PPT Shape 폰트변경 sub 함수 (변경할 Shape, Font 이름)
def C_Font (shp,Fname):
    for paragraph in shp.text_frame.paragraphs:
         for run in paragraph.runs:
            run.font.name = Fname
            
# PPT 폰트 통일 함수, Shape 형태에 따라 구분 동작 (파일명, 작업대상폴더, 결과저장폴더, 바꿀언어)
def AutoFont(file,Tpath,Rpath,Fname):

    prs = Presentation(Tpath + '/' + file)

    for slide in prs.slides:
        for shape in slide.shapes:
            #shape에 text frame이 있을 경우
             if shape.has_text_frame:
                 C_Font(shape,Fname)
            #shape가 Table 일 경우
             if shape.has_table:
                 for row in shape.table.rows:
                     for cell in row.cells:
                         C_Font(cell,Fname)
            #shape가 Group으로 묶여 있을 경우
             if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                 for shp in shape.shapes:
                     if shp.has_text_frame:
                         C_Font(shp,Fname)
    prs.save(Rpath + '/' + file)

# PPT 폰트 통일 함수2, Shape 형태에 따라 구분 동작 (파일명, 작업대상폴더, 결과저장폴더, 바꿀언어)
# PPT Application을 직접 실행하는 형태. 
def AutoFont2(file,Tpath,Rpath,FName):
    powerpoint = win32com.client.Dispatch('PowerPoint.Application')
    Sepa = '\\'
    Tpath = Tpath.replace("/","\\")
    Rpath = Rpath.replace("/","\\")
    ppt = powerpoint.Presentations.Open (Tpath + Sepa + file, WithWindow=False)

    for slide in ppt.Slides:
        for shape in slide.shapes:
            if shape.HasTextFrame == -1:
                shape.TextFrame.TextRange.Font.NameFarEast = FName
                shape.TextFrame.TextRange.Font.Name = FName
            if shape.HasTable == -1:
                for row in shape.Table.Rows:
                    for cell in row.cells:
                        cell.Shape.TextFrame.TextRange.Font.NameFarEast = FName
                        cell.Shape.TextFrame.TextRange.Font.Name = FName
            try:
                for GI in shape.GroupItems:
                    GI.TextFrame.TextRange.Font.NameFarEast = FName
                    GI.TextFrame.TextRange.Font.Name = FName
            except:
                pass

    ppt.SaveAs (Rpath + Sepa + file)
    ppt.Close ()
    
# PPT Replace Text Sub 함수 (Shape, 바뀔 keyword, 바꿀 keyword) 
def R_TXT (shp,B_W,A_W):
    for paragraph in shp.text_frame.paragraphs:
        for run in paragraph.runs:
            cur_text = run.text
            new_text = cur_text.replace(str(B_W), str(A_W))
            run.text = new_text

# PPT 키워드 변경, Shape 형태에 따라 구분 동작 (파일명, 작업대상폴더, 결과저장폴더, 바꿀언어)            
def PPT_KC (file,SPath,RPath,DPath):
    wb = openpyxl.load_workbook(DPath)
    ws = wb.active
    before_W = []
    after_W = []
    for cell in ws.rows:
        before_W.append(cell[0].value)
        after_W.append(cell[1].value)

    prs = Presentation(SPath + "/" + file)

    for slide in prs.slides:
        for shape in slide.shapes:
            for i in range(1, len(after_W)):
                # shape에 text frame이 있을 때
                if shape.has_text_frame:
                    R_TXT(shape,before_W[i],after_W[i])
                # shape가 Table일때    
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            R_TXT(cell,before_W[i],after_W[i])
                # shape가 Group으로 묶여 있을 때
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for shp in shape.shapes:
                        if shp.has_text_frame:
                            R_TXT(shp,before_W[i],after_W[i])
    prs.save(RPath + "/" + file)

